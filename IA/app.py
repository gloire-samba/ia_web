import os
import base64
import shutil
import uuid
from typing import Optional
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware

# Imports issus de tes fichiers et de ta stack technique

import time
import pandas as pd
import zipfile
import xml.etree.ElementTree as ET
import google.generativeai as genai
import subprocess
import io, warnings, csv, re
import docx2txt
import pymupdf4llm
import openpyxl
import xlsxwriter


from smolagents import CodeAgent, LiteLLMModel, tool, DuckDuckGoSearchTool, Tool
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_text_splitters import MarkdownTextSplitter
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor
from PIL import Image
from reportlab.pdfgen import canvas
try: from pdf2image import convert_from_path
except: pass

warnings.filterwarnings("ignore")

# --- CONFIGURATION SÉCURISÉE ---
# On récupère la clé depuis les secrets de l'hébergeur
api_key = os.environ.get("GOOGLE_API_KEY")

if not api_key:
    # Fallback pour tes tests locaux si la variable d'env n'est pas encore mise
    api_key = "TON_API_KEY_SI_BESOIN_LOCALEMENT" 

genai.configure(api_key=api_key)

# --- INITIALISATION DU MODÈLE (Le "Cerveau" de tes fichiers) ---
# Utilisation de Qwen ou Gemini selon ta configuration préférée dans le notebook
# --- INITIALISATION DU MODÈLE (Standard smolagents) ---
model = LiteLLMModel(
    model_id="gemini/gemini-2.5-flash", 
    api_key=api_key
)

# --- DÉFINITION DES OUTILS (Tools) ---
def convertir_tout_document(chemin_fichier):
    """Dispatcher Universel"""
    if not os.path.exists(chemin_fichier): return ""
    ext = os.path.splitext(chemin_fichier)[1].lower()
    texte = ""
    print(f"📂 Lecture ({ext}) : {os.path.basename(chemin_fichier)}")
    try:
        if ext == ".pdf": texte = pymupdf4llm.to_markdown(chemin_fichier)
        elif ext in [".jpg", ".png", ".jpeg"]:
            img = Image.open(chemin_fichier)
            texte = appel_gemini_securise("Décris cette image en détail.", img)
        elif ext == ".docx": texte = docx2txt.process(chemin_fichier)
        elif ext == ".pptx":
            prs = Presentation(chemin_fichier)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): texte += shape.text + "\n"
        elif ext in [".xlsx", ".xls"]:
            try: texte = pd.read_excel(chemin_fichier).fillna("").to_markdown(index=False)
            except: pass
        elif ext == ".txt":
            with open(chemin_fichier, 'r', encoding='utf-8', errors='ignore') as f: texte = f.read()
        elif ext == ".csv":
            try: texte = pd.read_csv(chemin_fichier).to_markdown(index=False)
            except: pass
    except Exception as e: return f"Erreur lecture globale: {e}"
    return texte

def appel_gemini_securise(prompt, image=None):
    """Appel Gemini Vision sécurisé avec RETRY et 2.5."""
    max_retries = 6; wait_time = 10
    try: model = genai.GenerativeModel("gemini-2.5-flash")
    except: model = genai.GenerativeModel("gemini-2.0-flash")

    for i in range(max_retries):
        try:
            if image: return model.generate_content([prompt, image]).text
            else: return model.generate_content(prompt).text
        except Exception as e:
            if "429" in str(e): time.sleep(wait_time); wait_time *= 2
            elif "404" in str(e): model = genai.GenerativeModel("gemini-2.5-flash")
            else: return f"Erreur Vision : {e}"
    return "Échec Quota."

embeddings = None

# TENTATIVE 1 : GOOGLE (Peut échouer en 404)
try:
    print("   👉 Tentative Google Embedding (004)...")
    test_emb = GoogleGenerativeAIEmbeddings(model="models/text-embedding-004")
    test_emb.embed_query("test")
    embeddings = test_emb
    print("✅ SUCCÈS : Mode Google activé.")
except Exception as e:
    print(f"   ❌ Échec Google (Erreur 404/API). Passage au Plan B.")

# TENTATIVE 2 : LOCAL (HUGGING FACE) - Infaillible
if embeddings is None:
    print("   👉 Activation du Plan B : Mémoire Locale (HuggingFace)...")
    try:
        # Modèle léger et très performant, gratuit, sans clé
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        print("✅ SUCCÈS : Mode Local (HuggingFace) activé. C'est plus robuste !")
    except Exception as e:
        raise Exception(f"🛑 CRITIQUE : Impossible de charger la mémoire locale. Erreur : {e}")

if embeddings is None:
    raise Exception("🛑 CRITIQUE : Échec total des mémoires.")


# --- RAG ---
vectorstore_global = None
def initialiser_rag(fichiers):
    global vectorstore_global
    text_data = ""
    for f in fichiers: text_data += convertir_tout_document(f) + "\n\n"
    if not text_data.strip(): return
    chunks = MarkdownTextSplitter(chunk_size=1000).split_text(text_data)
    vectorstore_global = FAISS.from_texts(chunks, embeddings)
    print("✅ Mémoire chargée.")

# ==============================================================================
# 5. OUTILS D'ACTION SPÉCIALISÉS (ARCHITECTURE V44)
# ==============================================================================

# --- A. OUTILS WORD ---

@tool
def createur_word(nom_fichier: str, contenu: str, style: str = "NORMAL", chemin_image: str = None) -> str:
    """
    Crée un NOUVEAU fichier Word (.docx).
    Args:
        nom_fichier: Nom du fichier à créer.
        contenu: Le texte à écrire.
        style: 'TITRE', 'GRAS', 'TAILLE_24', 'ROUGE', 'NORMAL'.
        chemin_image: Chemin d'une image à insérer (optionnel).
    """
    try:
        doc = Document()
        p = doc.add_paragraph()
        run = p.add_run(contenu)
        s = style.upper()
        if "GRAS" in s: run.bold = True
        if "ROUGE" in s: run.font.color.rgb = DocxRGBColor(255, 0, 0)
        if "TITRE" in s: p.style = 'Heading 1'
        if "TAILLE" in s:
            try: run.font.size = DocxPt(int(re.search(r'\d+', s).group()))
            except: pass
        if chemin_image: doc.add_picture(chemin_image, width=DocxInches(4))
        doc.save(nom_fichier)
        return "Word Créé."
    except Exception as e: return f"Erreur Création Word: {e}"

@tool
def modificateur_word(nom_fichier: str, texte_ancrage: str, texte_remplacement: str, action: str = "AJOUTER_FIN", chemin_image: str = None) -> str:
    """
    Modifie un fichier Word existant.
    Args:
        nom_fichier: Le fichier à modifier.
        texte_ancrage: Le texte à chercher (pour remplacement).
        texte_remplacement: Le nouveau texte (ou texte à ajouter).
        action: 'AJOUTER_FIN', 'REMPLACER'.
        chemin_image: Image à ajouter.
    """
    try:
        doc = Document(nom_fichier)
        if action == "REMPLACER":
            for p in doc.paragraphs:
                if texte_ancrage and texte_ancrage in p.text:
                    p.text = p.text.replace(texte_ancrage, texte_remplacement)
        else:
            doc.add_paragraph(texte_remplacement)
        if chemin_image: doc.add_picture(chemin_image, width=DocxInches(4))
        doc.save(nom_fichier)
        return "Word Modifié."
    except Exception as e: return f"Erreur Modif Word: {e}"
    
    
# --- B. OUTILS EXCEL ---

@tool
def createur_excel(nom_fichier: str, nom_feuille: str = "Donnees", donnees_initiales: str = "") -> str:
    """
    Crée un NOUVEAU fichier Excel (.xlsx).
    Args:
        nom_fichier: Nom du fichier.
        nom_feuille: Nom de la première feuille.
        donnees_initiales: (Optionnel) Données CSV brutes à écrire.
    """
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = nom_feuille
        if donnees_initiales:
            rows = donnees_initiales.split('\n')
            for r_idx, row in enumerate(rows, 1):
                cols = row.split(',')
                for c_idx, val in enumerate(cols, 1):
                    ws.cell(row=r_idx, column=c_idx, value=val.strip())
        wb.save(nom_fichier)
        return f"Excel créé : {nom_fichier}"
    except Exception as e: return f"Erreur Création Excel: {e}"

@tool
def modificateur_excel(nom_fichier: str, cible: str, valeur: str, action: str = "ECRIRE", style: str = None) -> str:
    """
    Modifie un Excel existant.
    Args:
        nom_fichier: Le fichier Excel.
        cible: La cellule (ex: 'B2') ou plage (ex: 'A1:C1').
        valeur: La donnée à écrire.
        action: 'ECRIRE', 'FUSION'.
        style: 'ROUGE', 'JAUNE', 'GRAS', 'TITRE'.
    """
    try:
        wb = openpyxl.load_workbook(nom_fichier)
        ws = wb.active
        val_clean = valeur
        try: val_clean = float(valeur.replace(' ','').replace('€','').replace(',','.'))
        except: pass

        if action == "ECRIRE": ws[cible] = val_clean
        elif action == "FUSION":
            ws.merge_cells(cible)
            top_left = cible.split(':')[0]
            ws[top_left] = val_clean
            ws[top_left].alignment = Alignment(horizontal='center', vertical='center')

        c = ws[cible.split(':')[0]]
        if style == "GRAS": c.font = Font(bold=True)
        if style == "ROUGE": c.font = Font(color="FF0000")
        if style == "JAUNE": c.fill = PatternFill(start_color="FFFF00", end_color="FFFF00", fill_type="solid")
        if style == "TITRE":
             c.font = Font(bold=True, color="FFFFFF", size=12)
             c.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        wb.save(nom_fichier)
        return "Excel Modifié."
    except Exception as e: return f"Erreur Modif Excel: {e}"
    
    
# --- C. OUTILS POWERPOINT (RECONSTRUCTION INTELLIGENTE) ---

@tool
def createur_ppt(nom_fichier: str, texte: str, chemin_image: str = None) -> str:
    """
    Crée un PowerPoint avec la logique "Flux Vertical" : Image -> Espace -> Texte.
    Gère automatiquement la création de nouvelles diapositives si le contenu déborde.
    Args:
        nom_fichier: Nom du fichier.
        texte: Texte long à insérer.
        chemin_image: Image à insérer (optionnel).
    """
    try:
        if os.path.exists(nom_fichier): os.remove(nom_fichier)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title: slide.shapes.title.text = "Présentation IA"
        cursor_y = Inches(1.5)

        # 1. IMAGE
        if chemin_image:
            img_height = Inches(3.5)
            slide.shapes.add_picture(chemin_image, Inches(1), cursor_y, height=img_height)
            cursor_y += img_height + Inches(0.2)

        # 2. TEXTE
        if texte:
            paragraphes = texte.split('\n')
            tb = slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
            tf = tb.text_frame
            tf.word_wrap = True

            for para in paragraphes:
                if not para.strip(): continue
                lines = max(1, len(para) / 90)
                height_needed = lines * 0.3

                if (cursor_y / Inches(1)) + height_needed > 7.5:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    cursor_y = Inches(1)
                    tb = slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
                    tf = tb.text_frame
                    tf.word_wrap = True

                p = tf.add_paragraph()
                p.text = para
                p.font.size = Pt(12)
                cursor_y += Inches(height_needed)

        prs.save(nom_fichier)
        return "PPT Créé (Flux continu respecté)."
    except Exception as e: return f"Erreur Création PPT: {e}"

@tool
def modificateur_ppt(nom_fichier: str, numero_slide: str, nouveau_texte: str, chemin_nouvelle_image: str = None) -> str:
    """
    Modifie une diapo en la RECONSTRUISANT proprement pour éviter les chevauchements.
    Si le texte déborde, crée de nouvelles slides.
    Args:
        nom_fichier: Le fichier PPTX.
        numero_slide: Numéro de la diapo.
        nouveau_texte: Le texte final à mettre.
        chemin_nouvelle_image: (Optionnel) Nouvelle image à mettre.
    """
    try:
        prs = Presentation(nom_fichier)
        idx = int(numero_slide) - 1
        slide = prs.slides[idx]
        # Nettoyage
        for shape in list(slide.shapes):
            if not shape == slide.shapes.title:
                sp = shape._element
                sp.getparent().remove(sp)

        # Reconstruction
        cursor_y = Inches(1.5)
        if chemin_nouvelle_image:
            img_height = Inches(3.5)
            slide.shapes.add_picture(chemin_nouvelle_image, Inches(1), cursor_y, height=img_height)
            cursor_y += img_height + Inches(0.2)

        if nouveau_texte:
            paragraphes = nouveau_texte.split('\n')
            tb = slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
            tf = tb.text_frame
            tf.word_wrap = True

            for para in paragraphes:
                if not para.strip(): continue
                lines = max(1, len(para) / 90)
                height_needed = lines * 0.3

                if (cursor_y / Inches(1)) + height_needed > 7.5:
                    new_slide = prs.slides.add_slide(prs.slide_layouts[5])
                    cursor_y = Inches(1)
                    tb = new_slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
                    tf = tb.text_frame
                    tf.word_wrap = True

                p = tf.add_paragraph()
                p.text = para
                p.font.size = Pt(12)
                cursor_y += Inches(height_needed)

        prs.save(nom_fichier)
        return "PPT Modifié (Slide reconstruite + Overflow géré)."
    except Exception as e: return f"Erreur Modif PPT: {e}"
    
# --- D. OUTILS DIVERS (TXT/CSV, PDF, VISION) ---

@tool
def editeur_texte_csv(nom_fichier: str, contenu: str, mode: str = "AJOUTER_FIN") -> str:
    """
    Crée ou modifie un fichier texte (.txt) ou CSV (.csv).

    Args:
        nom_fichier: Le nom du fichier.
        contenu: Le texte à écrire.
        mode: 'AJOUTER_FIN' (ajoute à la fin) ou 'ECRASER' (remplace tout).
    """
    try:
        m = 'w' if mode=="ECRASER" else 'a'
        with open(nom_fichier, m, encoding='utf-8') as f: f.write("\n"+contenu)
        return "Fichier Texte Modifié."
    except Exception as e: return f"Erreur TXT: {e}"

@tool
def convertisseur_pdf_vers_editable(chemin_source: str, type_sortie: str = "docx") -> str:
    """
    Convertit un fichier PDF vers un format éditable (docx ou xlsx).
    Args:
        chemin_source: Le chemin vers le fichier PDF d'origine.
        type_sortie: Le format de sortie souhaité, soit 'docx' soit 'xlsx'.
    """
    try:
        abs_in = os.path.abspath(chemin_source)
        dossier_sortie = os.path.dirname(abs_in)
        nom_sortie = os.path.splitext(chemin_source)[0] + "." + type_sortie
        
        args = ['libreoffice', '--headless', '--convert-to', type_sortie, abs_in, '--outdir', dossier_sortie]
        if type_sortie == 'xlsx': args.append('--infilter=CSV:44,34,76')
        subprocess.run(args, check=True, stdout=subprocess.DEVNULL)
        return f"Succès. Fichier converti en : {nom_sortie}."
    except Exception as e: return f"Erreur Conversion: {e}"

@tool
def convertisseur_editable_vers_pdf(chemin_source: str) -> str:
    """
    Convertit un fichier éditable (Word, Excel, PPT) vers le format PDF.
    Args:
        chemin_source: Le chemin vers le fichier source à convertir en PDF.
    """
    try:
        abs_in = os.path.abspath(chemin_source)
        dossier_sortie = os.path.dirname(abs_in)
        
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', abs_in, '--outdir', dossier_sortie], check=True, stdout=subprocess.DEVNULL)
        return "Succès. Reconverti en PDF."
    except Exception as e: return f"Erreur Conversion: {e}"
    

@tool
def outil_vision(chemin_image: str, question: str) -> str:
    """
    Analyse image avec Gemini.

    Args:
        chemin_image: Chemin de l'image.
        question: Question à poser.
    """
    try: return appel_gemini_securise(f"{question}", Image.open(chemin_image))
    except: return "Erreur Vision"

@tool
def outil_rag(question: str) -> str:
    """
    Cherche des informations UNIQUEMENT dans les documents fournis par l'utilisateur (RAG).
    NE L'UTILISE PAS pour des questions de culture générale.
    Args:
        question: La question.
    """
    global vectorstore_global
    if vectorstore_global is None: return "Aucun document fourni par l'utilisateur."
    try:
        docs = vectorstore_global.similarity_search(question, k=10)
        return "\n".join([d.page_content for d in docs])
    except: return "Erreur RAG"

@tool
def web_search(query: str) -> str:
    """
    Recherche Web (DuckDuckGo).

    Args:
        query: Mots-clés de la recherche.
    """
    try: 
        return DuckDuckGoSearchTool().run(query)
    except: 
        return "ERREUR_RESEAU. Ne cherche pas dans ta mémoire. Dis exactement ceci à l'utilisateur : 'Je n'ai pas pu vérifier cette information sur Internet à cause d'un blocage réseau temporaire. Pour éviter de vous donner des données fausses ou incomplètes, je préfère ne pas répondre de mémoire.'"
    
# La liste COMPLÈTE des imports dont l'IA a besoin pour la V44
imports_autorises = ["os", "pandas", "zipfile", "openpyxl", "pptx", "docx", "subprocess", "reportlab", "PIL", "csv", "pdf2image", "re", "time"]

agent = CodeAgent(
    tools=[outil_rag, outil_vision, createur_word, modificateur_word, createur_excel, modificateur_excel, createur_ppt, modificateur_ppt, editeur_texte_csv, convertisseur_pdf_vers_editable, convertisseur_editable_vers_pdf, web_search], 
    model=model,
    additional_authorized_imports=imports_autorises,
    max_steps=20
)


# On réinjecte les règles d'or (LE CERVEAU DE LA V44)
consigne = """
RÈGLES D'OR V44 :
1. ANALYSE LA DEMANDE : S'agit-il de manipuler un fichier (Word, Excel, PPT) ou est-ce une question de culture générale ?
2. CULTURE GÉNÉRALE : Si la question ne nécessite pas d'analyser un document fourni, réponds directement avec tes propres connaissances sans utiliser d'outil.
3. MODIF PPT : Utilise 'modificateur_ppt'. Il vide la slide et la recrée proprement avec le contenu final (Image + Texte + Overflow).
4. MODIF PDF : Convertis (docx/xlsx) -> Modifie -> Reconvertis.
5. HONNÊTETÉ : Si tu ne trouves pas une info, dis "NON TROUVÉE". Ne l'invente jamais.
"""
agent.prompt_templates["system_prompt"] = consigne + agent.prompt_templates["system_prompt"]

# --- MODÈLES DE DONNÉES API ---
app = FastAPI(title="IA Bureautique API")

# --- 2. AJOUT DU CORS POUR AUTORISER ANGULAR ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # En production pro, on mettrait ["http://localhost:4200"] au lieu de ["*"]
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class ChatRequest(BaseModel):
    prompt: str
    file_name: Optional[str] = None
    file_base64: Optional[str] = None

class ChatResponse(BaseModel):
    text: str
    output_file_name: Optional[str] = None
    output_file_base64: Optional[str] = None

# --- POINT D'ENTRÉE PRINCIPAL ---
# --- POINT D'ENTRÉE PRINCIPAL ---
@app.post("/api/chat")
async def process_request(request: ChatRequest):
    global vectorstore_global # On appelle la mémoire globale
    
    # Création d'un ID unique pour cette session de travail
    session_id = str(uuid.uuid4())
    work_dir = f"workspace_{session_id}"
    os.makedirs(work_dir, exist_ok=True)

    try:
        # 1. Gestion du fichier entrant et du RAG
        input_path = None
        if request.file_base64 and request.file_name:
            input_path = os.path.join(work_dir, request.file_name)
            with open(input_path, "wb") as f:
                f.write(base64.b64decode(request.file_base64))
            
            # CORRECTION : On force le RAG à lire le fichier entrant !
            initialiser_rag([input_path])
        else:
            # CORRECTION : S'il n'y a pas de fichier, on vide la mémoire 
            # pour ne pas utiliser le fichier de la requête précédente
            vectorstore_global = None 
            
        # 2. Exécution de l'IA
        # On demande à l'agent de travailler dans le dossier spécifique
        # 2. Exécution de l'IA
        # Construction de l'instruction de base
        instruction = (
            f"RÈGLE ABSOLUE : Ton espace de travail est le dossier '{work_dir}/'. "
            f"Tu DOIS OBLIGATOIREMENT préfixer tous les noms de fichiers que tu crées, modifies ou lis avec ce chemin exact.\n"
        )
        
        # CORRECTION : Si un fichier est fourni, on le dit explicitement à l'agent !
        if input_path:
            instruction += (
                f"\n⚠️ INFORMATION IMPORTANTE : L'utilisateur a joint un fichier nommé '{request.file_name}'. "
                f"Son chemin d'accès complet est '{input_path}'. Utilise tes outils sur ce fichier si la demande le nécessite.\n"
            )

        instruction += f"\nDemande de l'utilisateur : {request.prompt}"
        resultat_ia = agent.run(instruction)

        # 3. Détection de fichiers générés/modifiés
        out_name = None
        out_base64 = None
        
        files = os.listdir(work_dir)
        generated_files = [f for f in files if f != request.file_name]
        
        if generated_files:
            out_name = generated_files[0] 
            with open(os.path.join(work_dir, out_name), "rb") as f:
                out_base64 = base64.b64encode(f.read()).decode('utf-8')

        return ChatResponse(
            text=str(resultat_ia),
            output_file_name=out_name,
            output_file_base64=out_base64
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
    finally:
        # Nettoyage pour ne pas saturer le serveur
        if os.path.exists(work_dir):
            shutil.rmtree(work_dir)
        # Sécurité supplémentaire : on vide la mémoire à la fin
        vectorstore_global = None

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=7860)