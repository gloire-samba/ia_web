import os
import re
import subprocess
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from smolagents import tool

# ==============================================================================
# OUTILS WORD
# ==============================================================================

@tool
def createur_word(nom_fichier: str, contenu: str, style: str = "NORMAL", chemin_image: str = None) -> str:
    """
    Crée un NOUVEAU fichier Word (.docx).
    Args:
        nom_fichier: Nom du fichier à créer.
        contenu: Le texte à écrire.
        style: 'TITRE', 'GRAS', 'TAILLE_24', 'ROUGE', 'NORMAL'.
        chemin_image: Chemin d'une image à insérer (optionnel).
    """
    try:
        doc = Document()
        p = doc.add_paragraph()
        run = p.add_run(contenu)
        s = style.upper()
        if "GRAS" in s: run.bold = True
        if "ROUGE" in s: run.font.color.rgb = DocxRGBColor(255, 0, 0)
        if "TITRE" in s: p.style = 'Heading 1'
        if "TAILLE" in s:
            try: run.font.size = DocxPt(int(re.search(r'\d+', s).group()))
            except: pass
        if chemin_image: doc.add_picture(chemin_image, width=DocxInches(4))
        doc.save(nom_fichier)
        return "Word Créé."
    except Exception as e: return f"Erreur Création Word: {e}"

@tool
def modificateur_word(nom_fichier: str, texte_ancrage: str, texte_remplacement: str, action: str = "AJOUTER_FIN", chemin_image: str = None) -> str:
    """
    Modifie un fichier Word existant sans détruire la mise en page.
    Args:
        nom_fichier: Le fichier à modifier.
        texte_ancrage: Le texte à chercher (pour remplacement).
        texte_remplacement: Le nouveau texte (ou texte à ajouter).
        action: 'AJOUTER_FIN', 'REMPLACER'.
        chemin_image: Image à ajouter.
    """
    try:
        doc = Document(nom_fichier)
        if action == "REMPLACER":
            
            # Fonction interne pour remplacer délicatement sans écraser le style XML
            def remplacer_texte_delicatement(paragraphes):
                for p in paragraphes:
                    if texte_ancrage in p.text:
                        # On essaie d'abord de remplacer dans les "runs" (les blocs de style)
                        for run in p.runs:
                            if texte_ancrage in run.text:
                                run.text = run.text.replace(texte_ancrage, texte_remplacement)
                        # Si le texte est à cheval sur plusieurs runs, on force le remplacement brutal
                        if texte_ancrage in p.text:
                            p.text = p.text.replace(texte_ancrage, texte_remplacement)

            # 1. Fouiller les paragraphes classiques
            remplacer_texte_delicatement(doc.paragraphs)
            
            # 2. Fouiller les Tableaux (CRUCIAL pour les PDF convertis !)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        remplacer_texte_delicatement(cell.paragraphs)
        else:
            doc.add_paragraph(texte_remplacement)
            
        if chemin_image: doc.add_picture(chemin_image, width=DocxInches(4))
        doc.save(nom_fichier)
        return "Word Modifié avec succès."
    except Exception as e: 
        return f"Erreur Modif Word: {e}"
    
# ==============================================================================
# OUTILS EXCEL
# ==============================================================================

@tool
def createur_excel(nom_fichier: str, nom_feuille: str = "Donnees", donnees_initiales: str = "") -> str:
    """
    Crée un NOUVEAU fichier Excel (.xlsx).
    Args:
        nom_fichier: Nom du fichier.
        nom_feuille: Nom de la première feuille.
        donnees_initiales: (Optionnel) Données CSV brutes à écrire.
    """
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = nom_feuille
        if donnees_initiales:
            rows = donnees_initiales.split('\n')
            for r_idx, row in enumerate(rows, 1):
                cols = row.split(',')
                for c_idx, val in enumerate(cols, 1):
                    ws.cell(row=r_idx, column=c_idx, value=val.strip())
        wb.save(nom_fichier)
        return f"Excel créé : {nom_fichier}"
    except Exception as e: return f"Erreur Création Excel: {e}"

@tool
def modificateur_excel(nom_fichier: str, cible: str, valeur: str, action: str = "ECRIRE", style: str = None) -> str:
    """
    Modifie un Excel existant.
    Args:
        nom_fichier: Le fichier Excel.
        cible: La cellule (ex: 'B2') ou plage (ex: 'A1:C1').
        valeur: La donnée à écrire.
        action: 'ECRIRE', 'FUSION'.
        style: 'ROUGE', 'JAUNE', 'GRAS', 'TITRE'.
    """
    try:
        wb = openpyxl.load_workbook(nom_fichier)
        ws = wb.active
        val_clean = valeur
        try: val_clean = float(valeur.replace(' ','').replace('€','').replace(',','.'))
        except: pass

        if action == "ECRIRE": ws[cible] = val_clean
        elif action == "FUSION":
            ws.merge_cells(cible)
            top_left = cible.split(':')[0]
            ws[top_left] = val_clean
            ws[top_left].alignment = Alignment(horizontal='center', vertical='center')

        c = ws[cible.split(':')[0]]
        if style == "GRAS": c.font = Font(bold=True)
        if style == "ROUGE": c.font = Font(color="FF0000")
        if style == "JAUNE": c.fill = PatternFill(start_color="FFFF00", end_color="FFFF00", fill_type="solid")
        if style == "TITRE":
             c.font = Font(bold=True, color="FFFFFF", size=12)
             c.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        wb.save(nom_fichier)
        return "Excel Modifié."
    except Exception as e: return f"Erreur Modif Excel: {e}"
    
# ==============================================================================
# OUTILS POWERPOINT
# ==============================================================================

@tool
def createur_ppt(nom_fichier: str, texte: str, chemin_image: str = None) -> str:
    """
    Crée un PowerPoint avec la logique "Flux Vertical" : Image -> Espace -> Texte.
    Gère automatiquement la création de nouvelles diapositives si le contenu déborde.
    Args:
        nom_fichier: Nom du fichier.
        texte: Texte long à insérer.
        chemin_image: Image à insérer (optionnel).
    """
    try:
        if os.path.exists(nom_fichier): os.remove(nom_fichier)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title: slide.shapes.title.text = "Présentation IA"
        cursor_y = Inches(1.5)

        if chemin_image:
            img_height = Inches(3.5)
            slide.shapes.add_picture(chemin_image, Inches(1), cursor_y, height=img_height)
            cursor_y += img_height + Inches(0.2)

        if texte:
            paragraphes = texte.split('\n')
            tb = slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
            tf = tb.text_frame
            tf.word_wrap = True

            for para in paragraphes:
                if not para.strip(): continue
                lines = max(1, len(para) / 90)
                height_needed = lines * 0.3

                if (cursor_y / Inches(1)) + height_needed > 7.5:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    cursor_y = Inches(1)
                    tb = slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
                    tf = tb.text_frame
                    tf.word_wrap = True

                p = tf.add_paragraph()
                p.text = para
                p.font.size = Pt(12)
                cursor_y += Inches(height_needed)

        prs.save(nom_fichier)
        return "PPT Créé (Flux continu respecté)."
    except Exception as e: return f"Erreur Création PPT: {e}"

@tool
def modificateur_ppt(nom_fichier: str, numero_slide: str, nouveau_texte: str, chemin_nouvelle_image: str = None) -> str:
    """
    Modifie une diapo en la RECONSTRUISANT proprement pour éviter les chevauchements.
    Si le texte déborde, crée de nouvelles slides.
    Args:
        nom_fichier: Le fichier PPTX.
        numero_slide: Numéro de la diapo.
        nouveau_texte: Le texte final à mettre.
        chemin_nouvelle_image: (Optionnel) Nouvelle image à mettre.
    """
    try:
        prs = Presentation(nom_fichier)
        idx = int(numero_slide) - 1
        slide = prs.slides[idx]
        for shape in list(slide.shapes):
            if not shape == slide.shapes.title:
                sp = shape._element
                sp.getparent().remove(sp)

        cursor_y = Inches(1.5)
        if chemin_nouvelle_image:
            img_height = Inches(3.5)
            slide.shapes.add_picture(chemin_nouvelle_image, Inches(1), cursor_y, height=img_height)
            cursor_y += img_height + Inches(0.2)

        if nouveau_texte:
            paragraphes = nouveau_texte.split('\n')
            tb = slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
            tf = tb.text_frame
            tf.word_wrap = True

            for para in paragraphes:
                if not para.strip(): continue
                lines = max(1, len(para) / 90)
                height_needed = lines * 0.3

                if (cursor_y / Inches(1)) + height_needed > 7.5:
                    new_slide = prs.slides.add_slide(prs.slide_layouts[5])
                    cursor_y = Inches(1)
                    tb = new_slide.shapes.add_textbox(Inches(1), cursor_y, Inches(8), Inches(1))
                    tf = tb.text_frame
                    tf.word_wrap = True

                p = tf.add_paragraph()
                p.text = para
                p.font.size = Pt(12)
                cursor_y += Inches(height_needed)

        prs.save(nom_fichier)
        return "PPT Modifié (Slide reconstruite + Overflow géré)."
    except Exception as e: return f"Erreur Modif PPT: {e}"
    
# ==============================================================================
# OUTILS DIVERS (TXT, PDF)
# ==============================================================================

@tool
def editeur_texte_csv(nom_fichier: str, contenu: str, mode: str = "AJOUTER_FIN") -> str:
    """
    Crée ou modifie un fichier texte (.txt) ou CSV (.csv).

    Args:
        nom_fichier: Le nom du fichier.
        contenu: Le texte à écrire.
        mode: 'AJOUTER_FIN' (ajoute à la fin) ou 'ECRASER' (remplace tout).
    """
    try:
        m = 'w' if mode=="ECRASER" else 'a'
        with open(nom_fichier, m, encoding='utf-8') as f: f.write("\n"+contenu)
        return "Fichier Texte Modifié."
    except Exception as e: return f"Erreur TXT: {e}"

@tool
def convertisseur_pdf_vers_editable(chemin_source: str, type_sortie: str = "docx") -> str:
    """
    Convertit un fichier PDF vers un format éditable (docx ou xlsx).
    Args:
        chemin_source: Le chemin vers le fichier PDF d'origine.
        type_sortie: Le format de sortie souhaité, soit 'docx' soit 'xlsx'.
    """
    try:
        abs_in = os.path.abspath(chemin_source)
        dossier_sortie = os.path.dirname(abs_in)
        nom_sortie = os.path.splitext(chemin_source)[0] + "." + type_sortie
        
        if type_sortie.lower() == "docx":
            # 👉 CORRECTION : On utilise la librairie spécialisée pdf2docx
            from pdf2docx import Converter
            cv = Converter(abs_in)
            cv.convert(nom_sortie)
            cv.close()
        else:
            # LibreOffice reste indispensable pour le reste (xlsx, etc.)
            args = ['libreoffice', '--headless', '--convert-to', type_sortie, abs_in, '--outdir', dossier_sortie]
            if type_sortie == 'xlsx': args.append('--infilter=CSV:44,34,76')
            subprocess.run(args, check=True, stdout=subprocess.DEVNULL)
            
        # 👉 FILET DE SÉCURITÉ : On vérifie que le fichier a vraiment été créé !
        if not os.path.exists(nom_sortie):
            return f"Erreur : L'outil n'a pas pu créer le fichier {nom_sortie}."
            
        return f"Succès. Fichier converti en : {nom_sortie}."
    except Exception as e: 
        return f"Erreur Conversion: {e}"

@tool
def convertisseur_editable_vers_pdf(chemin_source: str) -> str:
    """
    Convertit un fichier éditable (Word, Excel, PPT) vers le format PDF.
    Args:
        chemin_source: Le chemin vers le fichier source à convertir en PDF.
    """
    try:
        abs_in = os.path.abspath(chemin_source)
        dossier_sortie = os.path.dirname(abs_in)
        base_name = os.path.splitext(os.path.basename(abs_in))[0]
        
        # Le nom que LibreOffice va donner par défaut
        pdf_genere_par_lo = os.path.join(dossier_sortie, base_name + ".pdf")
        # 👉 LE CORRECTIF : Le nouveau nom pour échapper au filtre "if f != file_name" de main.py
        pdf_final_renomme = os.path.join(dossier_sortie, base_name + "_modifie_IA.pdf")
        
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', abs_in, '--outdir', dossier_sortie], check=True, stdout=subprocess.DEVNULL)
        
        # On renomme le fichier immédiatement après la conversion
        if os.path.exists(pdf_genere_par_lo):
            os.rename(pdf_genere_par_lo, pdf_final_renomme)
            return f"Succès. Reconverti en PDF sous le nom : {pdf_final_renomme}"
            
        return "Erreur : Le fichier PDF n'a pas été généré."
    except Exception as e: 
        return f"Erreur Conversion: {e}"