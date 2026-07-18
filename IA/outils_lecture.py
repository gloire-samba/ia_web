import os
import time
import pandas as pd
import docx2txt
import pymupdf4llm
from PIL import Image
from pptx import Presentation
from smolagents import tool
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import MarkdownTextSplitter

# 👉 LE NOUVEL IMPORT OFFICIEL GOOGLE
from google import genai 

# ==============================================================================
# FONCTIONS DE LECTURE GLOBALE ET VISION
# ==============================================================================

def appel_gemini_securise(prompt, image=None):
    """Appel Gemini Vision sécurisé avec le NOUVEAU SDK (google-genai) et Gemini 2.5."""
    
    # 👉 CORRECTION MAJEURE : On empêche le serveur de dormir pendant 10 minutes !
    max_retries = 2  # Au lieu de 6
    wait_time = 2    # Au lieu de 10 secondes
    
    api_key = os.environ.get("GOOGLE_API_KEY")
    
    if not api_key:
        return "⚠️ Erreur : La clé GOOGLE_API_KEY est introuvable dans les secrets."
        
    try:
        client = genai.Client(api_key=api_key)
    except Exception as e:
        return f"Erreur Init Client : {e}"

    model_name = "gemini-2.5-flash"

    for i in range(max_retries):
        try:
            contents = [prompt, image] if image else prompt
            response = client.models.generate_content(
                model=model_name,
                contents=contents
            )
            return response.text
            
        except Exception as e:
            error_msg = str(e)
            # On ne fait qu'une toute petite pause de 2 secondes
            if "429" in error_msg or "503" in error_msg: 
                time.sleep(wait_time)
            elif "404" in error_msg: 
                model_name = "gemini-2.0-flash"
            else: 
                return f"Erreur Vision : {e}"
                
    # Si Google refuse toujours, on renvoie un texte simple au lieu de bloquer l'IA
    return "⚠️ Image illisible pour l'instant (Quota Google dépassé)."

def convertir_tout_document(chemin_fichier):
    """Dispatcher Universel"""
    if not os.path.exists(chemin_fichier): return ""
    ext = os.path.splitext(chemin_fichier)[1].lower()
    texte = ""
    print(f"📂 Lecture ({ext}) : {os.path.basename(chemin_fichier)}")
    try:
        if ext == ".pdf": texte = pymupdf4llm.to_markdown(chemin_fichier)
        elif ext in [".jpg", ".png", ".jpeg"]:
            img = Image.open(chemin_fichier)
            texte = appel_gemini_securise("Décris cette image en détail.", img)
        elif ext == ".docx": texte = docx2txt.process(chemin_fichier)
        elif ext == ".pptx":
            prs = Presentation(chemin_fichier)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): texte += shape.text + "\n"
        elif ext in [".xlsx", ".xls"]:
            try: texte = pd.read_excel(chemin_fichier).fillna("").to_markdown(index=False)
            except: pass
        # 👉 CORRECTION 2 : Le RAG lit instantanément tous ces langages comme du texte !
        elif ext in [".txt", ".py", ".js", ".ts", ".html", ".css", ".java", ".md", ".json", ".yml", ".sh"]:
            with open(chemin_fichier, 'r', encoding='utf-8', errors='ignore') as f: texte = f.read()
        elif ext == ".csv":
            try: texte = pd.read_csv(chemin_fichier).to_markdown(index=False)
            except: pass
    except Exception as e:
        return (
            f"ERREUR_CRITIQUE_OUTIL : {str(e)}\n"
            "INSTRUCTION STRICTE POUR L'AGENT : Tu ne dois plus essayer d'utiliser cet outil. "
            "Tu dois OBLIGATOIREMENT contourner le problème en répondant à la demande de l'utilisateur "
            "uniquement grâce à tes connaissances internes. "
            "Commence ta réponse finale par CETTE phrase exacte :\n"
            "\"⚠️ **L'outil demandé ou le réseau est indisponible. Je vous réponds avec mes connaissances personnelles. Veuillez garder à l'esprit que je suis un modèle gratuit et qu'il est possible que je me trompe.** ⚠️\""
        )
    return texte

@tool
def outil_vision(chemin_image: str, question: str) -> str:
    """
    Analyse image avec Gemini.
    Args:
        chemin_image: Chemin de l'image.
        question: Question à poser.
    """
    try: return appel_gemini_securise(f"{question}", Image.open(chemin_image))
    except Exception as e:
        return (
            f"ERREUR_CRITIQUE_OUTIL : {str(e)}\n"
            "INSTRUCTION STRICTE POUR L'AGENT : Tu ne dois plus essayer d'utiliser cet outil. "
            "Tu dois OBLIGATOIREMENT contourner le problème en répondant à la demande de l'utilisateur "
            "uniquement grâce à tes connaissances internes. "
            "Commence ta réponse finale par CETTE phrase exacte :\n"
            "\"⚠️ **L'outil demandé ou le réseau est indisponible. Je vous réponds avec mes connaissances personnelles. Veuillez garder à l'esprit que je suis un modèle gratuit et qu'il est possible que je me trompe.** ⚠️\""
        )

# ==============================================================================
# GESTION DU RAG (MÉMOIRE TEMPORAIRE DE L'AGENT)
# ==============================================================================

vectorstore_global = None

def initialiser_rag(fichiers, embeddings):
    """Initialise le RAG avec les documents fournis et les embeddings injectés"""
    global vectorstore_global
    text_data = ""
    for f in fichiers: text_data += convertir_tout_document(f) + "\n\n"
    if not text_data.strip(): return
    chunks = MarkdownTextSplitter(chunk_size=1000).split_text(text_data)
    vectorstore_global = FAISS.from_texts(chunks, embeddings)
    print("✅ Mémoire chargée.")

def vider_memoire_rag():
    """Vide la mémoire RAG entre deux requêtes API"""
    global vectorstore_global
    vectorstore_global = None

@tool
def outil_rag(question: str) -> str:
    """
    Cherche des informations UNIQUEMENT dans les documents fournis par l'utilisateur (RAG).
    NE L'UTILISE PAS pour des questions de culture générale.
    Args:
        question: La question.
    """
    global vectorstore_global
    if vectorstore_global is None: return "Aucun document fourni par l'utilisateur."
    try:
        docs = vectorstore_global.similarity_search(question, k=10)
        return "\n".join([d.page_content for d in docs])
    except Exception as e:
        return (
            f"ERREUR_CRITIQUE_OUTIL : {str(e)}\n"
            "INSTRUCTION STRICTE POUR L'AGENT : Tu ne dois plus essayer d'utiliser cet outil. "
            "Tu dois OBLIGATOIREMENT contourner le problème en répondant à la demande de l'utilisateur "
            "uniquement grâce à tes connaissances internes. "
            "Commence ta réponse finale par CETTE phrase exacte :\n"
            "\"⚠️ **L'outil demandé ou le réseau est indisponible. Je vous réponds avec mes connaissances personnelles. Veuillez garder à l'esprit que je suis un modèle gratuit et qu'il est possible que je me trompe.** ⚠️\""
        )
